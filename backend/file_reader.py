"""Universal file content extraction for uploaded documents.

Supports PDF, DOCX, XLSX, PPTX, CSV/TSV and a wide range of plain-text
formats (TXT/MD/JSON/HTML/XML/YAML/LOG and code files). Images and audio
are handled by the vision/voice services elsewhere.
"""
import csv
import io
import logging

logger = logging.getLogger("ai_os.file_reader")

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".json", ".jsonl", ".html", ".htm",
    ".xml", ".yaml", ".yml", ".log", ".ini", ".cfg", ".conf",
    ".py", ".js", ".ts", ".tsx", ".sql", ".sh", ".css", ".env",
    ".csv", ".tsv",
}

BINARY_HINT = "[Binary or unsupported file"


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1", "cp1252"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")


def _looks_like_text(data: bytes) -> bool:
    sample = data[:4096]
    if not sample:
        return True
    printable = sum(1 for c in sample if c in (9, 10, 13) or 32 <= c < 127 or c > 160)
    return printable / len(sample) > 0.85


def extract_text(filename: str, data: bytes) -> str:
    """Return human-readable text extracted from an uploaded file."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == ".pdf":
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if ext == ".docx":
        import docx
        document = docx.Document(io.BytesIO(data))
        parts = [p.text for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                parts.append(" | ".join(cells))
        return "\n".join(parts)

    if ext == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        out = []
        try:
            for ws in wb.worksheets:
                out.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    vals = ["" if v is None else str(v).replace("\n", " ") for v in row]
                    if any(vals):
                        out.append(" | ".join(vals))
        finally:
            wb.close()
        return "\n".join(out)

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    out.append(shape.text)
        return "\n".join(out)

    if ext in (".csv", ".tsv"):
        delimiter = "\t" if ext == ".tsv" else ","
        text = _decode_text(data)
        try:
            rows = list(csv.reader(io.StringIO(text), delimiter=delimiter))
            return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)
        except Exception:
            return text

    if ext in TEXT_EXTENSIONS:
        return _decode_text(data)

    if _looks_like_text(data):
        return _decode_text(data)

    return f"{BINARY_HINT}: {filename or 'file'}]"
